from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

DARK_BG = RGBColor(15, 23, 42)
SURFACE = RGBColor(17, 29, 53)
CYAN = RGBColor(14, 165, 233)
TEXT = RGBColor(232, 239, 246)
TEXT_SEC = RGBColor(139, 149, 176)
BORDER = RGBColor(30, 42, 69)

def blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide

# Slide 1: Cover
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
tf = title_box.text_frame
tf.text = "VitaTwin AI"
p = tf.paragraphs[0]
p.font.size = Pt(66)
p.font.bold = True
p.font.color.rgb = TEXT
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.6))
tf = sub_box.text_frame
tf.text = "Advanced Health Monitoring Platform"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = CYAN
p.alignment = PP_ALIGN.CENTER

# Slide 2: Login
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 1: Authentication"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(4.2))
shape.fill.solid()
shape.fill.fore_color.rgb = SURFACE
shape.line.color.rgb = BORDER

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Login & Sign Up Page\n\nSign In tab - Login with email and password\nSign Up tab - Create new account\nEmail validation and password requirements\nError handling with clear messages\nDemo mode - Use any email/password (min 6 chars)"
for p in tf.paragraphs:
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT

# Slide 3: Profile
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 2: User Profile Dashboard"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Profile Card:\nAvatar with name, Email, Premium status, Member date, Logout button\n\nHealth Stats:\nHealth Score: 87, Risk Level: Low, Age: 24 years, Height: 5 ft 10 in, Weight: 72 kg, Blood Type: O+"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 4: Dashboard
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 3: Health Dashboard"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "4 Key Metrics: Health Score, Age Estimate, Risk Level, Conditions\nDigital Twin Status: Visual representation with organ health\nOrgan Health: Heart 98%, Brain 96%, Lungs 94%, Immune 90%\nToday Activity: Steps 8432, Calories 1234, Sleep 7h 45m, Stress Low\nWeekly Trend: Line chart showing heart rate and sleep patterns"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 5: Vitals
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 4: Vitals Tracker"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Track and Log:\nWeight, Blood Pressure (Systolic/Diastolic), Heart Rate, Mood (1-10 scale)\n\nData Display:\nHistorical tracking, Trend visualization, Health insights, Progress analysis, Chart display"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 6: Progress
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 5: Progress Dashboard"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "7-Day Health Trends:\nWeight trend line chart\nBlood pressure trends (systolic and diastolic)\nHeart rate patterns\nMood tracking over 7 days\nAverage metrics calculation\nWeek-over-week comparison"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 7: Appointments
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 6: Doctor Appointments Calendar"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Appointment Management:\nAdd new appointments\nDoctor name and specialty selection\nDate and time scheduling\nLocation and notes\nUpcoming appointments list\nDelete and edit appointments\nCalendar view"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 8: Treatment Simulator
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 7: Treatment Simulator"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "AI-Powered Treatment Prediction:\nEffectiveness prediction (0-100%)\nLikely side effects list\nDrug interaction severity assessment\nEstimated onset time\nCost analysis (generic vs brand)\nDosage calculator\nPersonalized recommendations"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 9: Symptoms
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 8: Symptom Tracker"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Log and Monitor Symptoms:\nSymptom name and description\nDuration tracking\nSeverity rating (1-10 scale)\nAssociated conditions\nWhen to see doctor alert\nSymptom history\nPattern detection"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 10: Drug Checker
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 9: Drug Interaction Checker"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Check Drug Safety:\nCheck drug interactions\nContraindication alerts\nAllergy detection\nInteraction severity (None/Mild/Moderate/Severe)\nAlternative medication suggestions\nDetailed risk explanation\nWhen to contact doctor"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 11: Digital Twin
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 10: Digital Twin Simulator"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Simulate Health Scenarios:\nWeight Loss - Fitness Plan\nHigh Stress - Work Pressure\nBetter Sleep - 8+ Hours\nNew Medication - Treatment effects\n\nPredicted Changes: Energy levels, Immunity, Heart health, Mental clarity"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 12: AI Diagnosis
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 11: AI Diagnosis Engine"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "AI-Powered Analysis:\nSymptom Analysis (Completed)\nAI Processing (Completed)\nDeep Scan (In Progress)\nCross-Check (Pending)\n\nTop Findings: Viral Pneumonia 86%, Bronchitis 45%, Cold/Flu 23%"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 13: Reminders
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 12: Medication Reminders"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Never Miss Medication:\nAdd medications\nSet reminder times\nDosage tracking\nBrowser notifications\nVoice reminders\nAlert 5 minutes before\nText-to-speech support\nUpcoming list and status"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 14: Settings
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Step 13: User Settings & Controls"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(3.8))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Personalization:\nNotification preferences\nVoice reminder toggle\nPrivacy settings\n\nData Management:\nExport health data\nView statistics\nDownload as JSON\nClear data option"
for p in tf.paragraphs:
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT

# Slide 15: Features Summary
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "Complete Feature Set"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = TEXT

text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.6), Inches(4.2))
tf = text_box.text_frame
tf.word_wrap = True
tf.text = "Professional Authentication\nReal-time Health Dashboard\nVitals and Weight Tracking\nAI-Powered Diagnostics\nTreatment Simulator\nDrug Interaction Checker\nDigital Twin Simulations\nAppointment Calendar\nMedication Reminders\nProgress Analytics\nSymptom Tracking\nUser Profile Management"
for p in tf.paragraphs:
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT

# Slide 16: Final
slide = blank_slide()
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
tf = title_box.text_frame
tf.text = "Ready to Use"
p = tf.paragraphs[0]
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = CYAN
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
tf = sub_box.text_frame
tf.text = "VitaTwin AI - Complete Health Platform"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.color.rgb = TEXT
p.alignment = PP_ALIGN.CENTER

prs.save('/sessions/wizardly-tender-euler/mnt/vitatwin-complete/VitaTwin_AI_Complete_Walkthrough.pptx')
print("Presentation created: VitaTwin_AI_Complete_Walkthrough.pptx")
