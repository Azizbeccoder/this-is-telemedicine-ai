from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Colors
DARK_BG = RGBColor(15, 23, 42)
SURFACE = RGBColor(17, 29, 53)
CYAN = RGBColor(14, 165, 233)
PURPLE = RGBColor(139, 92, 246)
PINK = RGBColor(236, 72, 153)
TEXT = RGBColor(232, 239, 246)
TEXT_SEC = RGBColor(139, 149, 176)
TEXT_MUTED = RGBColor(90, 100, 120)
BORDER = RGBColor(30, 42, 69)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(66)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEXT
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.6))
    sub_frame = sub_box.text_frame
    sub_frame.text = subtitle
    sub_frame.paragraphs[0].font.size = Pt(28)
    sub_frame.paragraphs[0].font.color.rgb = CYAN
    sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_left, content_right=None):
    """Add a content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEXT

    # Content box
    if content_right is None:
        # Full width content
        shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(9), Inches(4.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SURFACE
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(1)

        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(3.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = content_left

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = TEXT
            paragraph.space_before = Pt(4)
            paragraph.space_after = Pt(4)
    else:
        # Two column layout
        # Left box
        shape_left = slide.shapes.add_shape(1, Inches(0.5), Inches(1), Inches(4.3), Inches(4.2))
        shape_left.fill.solid()
        shape_left.fill.fore_color.rgb = SURFACE
        shape_left.line.color.rgb = BORDER
        shape_left.line.width = Pt(1)

        text_box_left = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(3.9), Inches(3.9))
        text_frame_left = text_box_left.text_frame
        text_frame_left.word_wrap = True
        text_frame_left.text = content_left

        for paragraph in text_frame_left.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = TEXT

        # Right box
        shape_right = slide.shapes.add_shape(1, Inches(5.2), Inches(1), Inches(4.3), Inches(4.2))
        shape_right.fill.solid()
        shape_right.fill.fore_color.rgb = SURFACE
        shape_right.line.color.rgb = BORDER
        shape_right.line.width = Pt(1)

        text_box_right = slide.shapes.add_textbox(Inches(5.4), Inches(1.2), Inches(3.9), Inches(3.9))
        text_frame_right = text_box_right.text_frame
        text_frame_right.word_wrap = True
        text_frame_right.text = content_right

        for paragraph in text_frame_right.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = TEXT

# Slide 1: Title
add_title_slide(prs, 'VitaTwin AI', 'Advanced Health Monitoring Platform')

# Slide 2: Login/SignUp
add_content_slide(prs, 'Step 1: Authentication',
'''Login & Sign Up Page

✓ Sign In tab - Login with email & password
✓ Sign Up tab - Create new account
✓ Email validation & password requirements
✓ Error handling with clear messages
✓ Demo mode - Use any email/password (min 6 chars)
✓ Beautiful gradient design
✓ Professional form validation''')

# Slide 3: User Profile
add_content_slide(prs, 'Step 2: User Profile Dashboard',
'''Profile Card

✓ Avatar with name
✓ Email address
✓ Premium status
✓ Member since date
✓ Logout button
✓ Quick access menu
✓ Professional styling''',
'''Health Stats & Info

Health Score: 87
Risk Level: Low
Age: 24 years
Height: 5'10"
Weight: 72 kg
Blood Type: O+
Total Checkups: 5''')

# Slide 4: Dashboard
add_content_slide(prs, 'Step 3: Health Dashboard',
'''Main Dashboard Overview

✓ 4 Key Metrics: Health Score, Age, Risk Level, Conditions
✓ Digital Twin Status: Visual representation
✓ Organ Details: Heart (98%), Brain (96%), Lungs (94%), Immune (90%)
✓ Today's Activity: Steps, Calories, Sleep, Stress
✓ Weekly Health Trend: Line chart with historical data
✓ Professional metrics display
✓ Real-time health insights''')

# Slide 5: Vitals Tracker
add_content_slide(prs, 'Step 4: Vitals Tracker',
'''Track Your Health Metrics

Log & Monitor:
✓ Weight
✓ Blood Pressure (Systolic/Diastolic)
✓ Heart Rate
✓ Mood (1-10 scale)
✓ Historical data
✓ Trend visualization''',
'''Data Display:

✓ Complete history
✓ Trend visualization
✓ Health insights
✓ Progress tracking
✓ Chart analysis
✓ Health recommendations''')

# Slide 6: Progress Dashboard
add_content_slide(prs, 'Step 5: Progress Dashboard',
'''7-Day Health Trends Analysis

Track weekly progress:
✓ Weight trend (line chart)
✓ Blood pressure trends (systolic & diastolic)
✓ Heart rate patterns
✓ Mood tracking over 7 days
✓ Average metrics calculation
✓ Week-over-week comparison
✓ Detailed analytics''')

# Slide 7: Doctor Appointments
add_content_slide(prs, 'Step 6: Doctor Appointments Calendar',
'''Schedule & Manage Appointments

✓ Add new appointments
✓ Doctor name & specialty selection
✓ Date & time scheduling
✓ Location & notes
✓ Upcoming appointments list
✓ Delete/edit appointments
✓ Calendar view''')

# Slide 8: Treatment Simulator
add_content_slide(prs, 'Step 7: Treatment Simulator',
'''AI-Powered Treatment Prediction

Select treatment & get:
✓ Effectiveness prediction (0-100%)
✓ Likely side effects list
✓ Drug interaction severity
✓ Estimated onset time
✓ Cost analysis (generic vs brand)
✓ Dosage calculator
✓ Personalized recommendations''')

# Slide 9: Symptom Tracker
add_content_slide(prs, 'Step 8: Symptom Tracker',
'''Log & Monitor Your Symptoms

Track symptoms with:
✓ Symptom name & description
✓ Duration tracking (days, hours)
✓ Severity rating (1-10 scale)
✓ Associated conditions
✓ When to see doctor alert
✓ Symptom history
✓ Pattern detection''')

# Slide 10: Drug Interaction Checker
add_content_slide(prs, 'Step 9: Drug Interaction Checker',
'''Check Drug Safety & Interactions

Comprehensive analysis includes:
✓ Check drug interactions
✓ Contraindication alerts
✓ Allergy detection
✓ Interaction severity (None/Mild/Moderate/Severe)
✓ Alternative medication suggestions
✓ Detailed explanation of risks
✓ When to contact doctor''')

# Slide 11: Digital Twin
add_content_slide(prs, 'Step 10: Digital Twin Simulator',
'''Simulate Health Scenarios

Test different scenarios:
✓ Weight Loss - Fitness Plan
✓ High Stress - Work Pressure
✓ Better Sleep - 8+ Hours
✓ New Medication - Treatment effects
✓ Lifestyle changes
✓ Health predictions''',
'''See predicted changes:

✓ Energy levels
✓ Immunity boost
✓ Heart health
✓ Mental clarity
✓ Overall wellness
✓ Risk reduction''')

# Slide 12: AI Diagnosis
add_content_slide(prs, 'Step 11: AI Diagnosis Engine',
'''AI-Powered Symptom Analysis

Complete diagnosis workflow:
✓ Symptom Analysis (Completed)
✓ AI Processing (Completed)
✓ Deep Scan (In Progress)
✓ Cross-Check (Pending)''',
'''Top Findings:

• Viral Pneumonia: 86%
• Bronchitis: 45%
• Cold/Flu: 23%

+ View full report''')

# Slide 13: Medication Reminders
add_content_slide(prs, 'Step 12: Medication Reminders',
'''Never Miss a Medication

Features:
✓ Add medications
✓ Set reminder times
✓ Dosage tracking
✓ Browser notifications
✓ Voice reminders''',
'''Notifications:

✓ Alert 5 min before
✓ Text-to-speech
✓ Visual indicators
✓ Upcoming list
✓ Mark as taken''')

# Slide 14: User Settings
add_content_slide(prs, 'Step 13: User Settings & Controls',
'''Personalize Your Experience

Control:
✓ Notification preferences
✓ Voice reminder toggle
✓ Privacy settings
✓ Data preferences''',
'''Data Management:

✓ Export health data
✓ View statistics
✓ Download as JSON
✓ Clear data option''')

# Slide 15: Features Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_BG

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
title_frame = title_box.text_frame
title_frame.text = 'Complete Feature Set'
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = TEXT

features = '''✓ Professional Authentication System
✓ Real-time Health Dashboard
✓ Vitals & Weight Tracking
✓ AI-Powered Diagnostics
✓ Treatment Simulator with Predictions
✓ Drug Interaction Checker
✓ Digital Twin Simulations
✓ Appointment Calendar
✓ Medication Reminders with Notifications
✓ Progress Analytics with Charts
✓ Symptom Tracking
✓ User Profile Management'''

text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.6), Inches(4.2))
text_frame = text_box.text_frame
text_frame.word_wrap = True
text_frame.text = features

for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(14)
    paragraph.font.color.rgb = TEXT
    paragraph.space_before = Pt(6)

# Slide 16: Final Slide
add_title_slide(prs, 'Ready to Use', 'VitaTwin AI - Complete Health Platform')

# Save presentation
output_path = '/sessions/wizardly-tender-euler/mnt/vitatwin-complete/VitaTwin_AI_Complete_Walkthrough.pptx'
prs.save(output_path)
print('Presentation created successfully!')
print('File: VitaTwin_AI_Complete_Walkthrough.pptx')
print('Total slides: 16')
